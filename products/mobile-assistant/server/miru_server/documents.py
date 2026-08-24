"""文档的本机解析与可选页渲染。

文字、表格、公式先留在本机提取；只有预览 PNG 才会进入 Vision 请求。
依赖可选，缺失时返回明确的可操作错误而不会损坏原文件。
"""
from __future__ import annotations

import csv
import json
import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class ExtractionResult:
    text: str = ""
    previews: list[str] = field(default_factory=list)
    error: str = ""


def extract(path: str | Path, kind: str, max_preview_pages: int = 10) -> ExtractionResult:
    path = Path(path)
    try:
        if kind == "text":
            return ExtractionResult(text=_read_text(path))
        if path.suffix.lower() == ".csv":
            return ExtractionResult(text=_read_csv(path))
        if path.suffix.lower() in {".xlsx", ".xls"}:
            return _read_workbook(path)
        if path.suffix.lower() == ".docx":
            return _read_docx(path)
        if path.suffix.lower() == ".pptx":
            return _read_pptx(path)
        if path.suffix.lower() == ".pdf":
            return _read_pdf(path, max_preview_pages)
        return ExtractionResult(error="暂不支持该文件格式")
    except Exception as e:
        return ExtractionResult(error=f"本机解析失败：{type(e).__name__}: {str(e)[:160]}")


def _read_text(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return data.decode(encoding)[:100_000]
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")[:100_000]


def _read_csv(path: Path) -> str:
    text = _read_text(path)
    rows = list(csv.reader(text.splitlines()))[:500]
    if not rows:
        return "空 CSV 文件"
    width = max((len(row) for row in rows), default=0)
    lines = [f"CSV：{len(rows)} 行（已读取前 500 行），最多 {width} 列", ""]
    for row in rows[:80]:
        lines.append(" | ".join(cell.strip().replace("\n", " ")[:160] for cell in row))
    return "\n".join(lines)


def _read_workbook(path: Path) -> ExtractionResult:
    try:
        import openpyxl
    except ImportError:
        return ExtractionResult(error="解析 Excel 需要 openpyxl；请在服务器安装 document 依赖")
    # 部分 App（包括 Cashew）导出的 xlsx 会把 worksheet dimension 错写成
    # A1:A1。只读模式因此把真实的 1194x11 明细误报成 1x1；空白封面表则
    # 可能得到 None。对可疑维度强制重扫 XML，仍保持只读模式以控制大文件内存。
    workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
    for sheet in workbook.worksheets:
        if (
            sheet.max_row is None
            or sheet.max_column is None
            or (sheet.max_row <= 1 and sheet.max_column <= 1)
        ):
            sheet.reset_dimensions()
            try:
                sheet.calculate_dimension(force=True)
            except UnboundLocalError:
                # openpyxl 对完全空白且无 dimension 的工作表会触发这个已知边界；
                # 保留 None，后续按空工作表处理即可。
                pass

    workbook_sheet_count = len(workbook.sheetnames)
    sheet_sections: list[tuple[int, int, str]] = []
    sheet_index: list[str] = []
    for index, name in enumerate(workbook.sheetnames[:20]):
        sheet = workbook[name]
        max_row = int(sheet.max_row or 0)
        max_column = int(sheet.max_column or 0)
        sheet_index.append(f"{name}（{max_row} 行 × {max_column} 列）")
        lines = [f"\n## 工作表：{name}（最大行 {max_row}，最大列 {max_column}）"]
        non_empty_rows = 0
        row_limit = min(max_row, 5000)
        column_limit = min(max_column, 60)
        if row_limit and column_limit:
            for row in sheet.iter_rows(
                min_row=1,
                max_row=row_limit,
                min_col=1,
                max_col=column_limit,
                values_only=True,
            ):
                values = [
                    "" if value is None else str(value).replace("\n", " ")[:240]
                    for value in row
                ]
                if any(values):
                    non_empty_rows += 1
                    lines.append(" | ".join(values))
        if non_empty_rows == 0:
            lines.append("（空工作表，已跳过）")
        if max_row > row_limit or max_column > column_limit:
            lines.append(
                f"（工作表过大，本机提取上限为 {row_limit} 行 × {column_limit} 列；"
                "其余内容未进入本轮模型上下文）"
            )
        sheet_sections.append((non_empty_rows, index, "\n".join(lines)))

    workbook.close()
    # 先放紧凑的汇总表，再放大体量明细表，防止明细表耗尽上下文后让后面的
    # 月度/分类汇总完全不可见。空表只在索引和末尾标明。
    sheet_sections.sort(key=lambda item: (item[0] == 0, item[0], item[1]))
    blocks = [
        f"Excel 文件：{path.name}；工作表 {workbook_sheet_count} 个",
        "工作表索引：" + "；".join(sheet_index),
        "以下按数据量从小到大展开，确保所有工作表都被覆盖。",
    ]
    for _, _, section in sheet_sections:
        remaining = 100_000 - len("\n".join(blocks))
        if remaining <= 0:
            break
        if len(section) > remaining:
            blocks.append(section[: max(remaining - 80, 0)] + "\n（正文达到本机提取上限，已截断）")
            break
        blocks.append(section)
    return ExtractionResult(text="\n".join(blocks)[:100_000])


def _read_docx(path: Path) -> ExtractionResult:
    try:
        from docx import Document
    except ImportError:
        return ExtractionResult(error="解析 Word 需要 python-docx；请在服务器安装 document 依赖")
    document = Document(path)
    blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for index, table in enumerate(document.tables, 1):
        blocks.append(f"\n## 表格 {index}")
        for row in table.rows:
            blocks.append(" | ".join(cell.text.replace("\n", " ")[:240] for cell in row.cells))
    return ExtractionResult(text="\n".join(blocks)[:100_000])


def _read_pptx(path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(error="解析 PPT 需要 python-pptx；请在服务器安装 document 依赖")
    presentation = Presentation(path)
    blocks: list[str] = [f"PPT：{len(presentation.slides)} 页"]
    for number, slide in enumerate(presentation.slides, 1):
        texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        blocks.append(f"\n## 第 {number} 页\n" + "\n".join(texts))
    return ExtractionResult(text="\n".join(blocks)[:100_000])


def _read_pdf(path: Path, max_pages: int) -> ExtractionResult:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ExtractionResult(error="解析 PDF/生成页面预览需要 PyMuPDF；请在服务器安装 document 依赖")
    document = fitz.open(path)
    blocks = [f"PDF：{document.page_count} 页"]
    preview_dir = path.parent / "previews"
    preview_dir.mkdir(exist_ok=True)
    previews: list[str] = []
    for index, page in enumerate(document):
        if index < 200:
            blocks.append(f"\n## 第 {index + 1} 页\n{page.get_text('text').strip()}")
        if index < max(max_pages, 0):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            target = preview_dir / f"{path.stem}-page-{index + 1}.png"
            pixmap.save(target)
            previews.append(str(target))
    document.close()
    return ExtractionResult(text="\n".join(blocks)[:100_000], previews=previews)
